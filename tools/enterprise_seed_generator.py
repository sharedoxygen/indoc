#!/usr/bin/env python3
"""
Enterprise-Grade Seed Data Generator for inDoc

Generates thousands of unique, professional business documents with:
- Realistic business content across multiple domains
- Comprehensive titles and descriptions
- Proper Elasticsearch and Qdrant indexing
- Equitable access distribution among users
- Real-world business scenarios
"""

import asyncio
import sys
import os
import logging
import hashlib
import random
import uuid
from datetime import datetime, timedelta
from pathlib import Path
from typing import Dict, Any, List
from io import BytesIO

# Add both app and backend directory to the Python path
sys.path.insert(0, str(Path(__file__).parent.parent / "backend"))
sys.path.insert(0, str(Path(__file__).parent.parent))

from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select, delete, func as sqlfunc
from faker import Faker

# Use backend imports (which have the correct paths)
import sys
sys.path = [p for p in sys.path if 'backend' not in p or p.endswith('backend')]  # Prioritize backend

from app.db.session import AsyncSessionLocal
from app.models.user import User, UserRole
from app.models.document import Document
from app.models.audit import AuditLog
from app.core.security import get_password_hash

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

fake = Faker()


class EnterpriseDocumentGenerator:
    """Generate enterprise-grade realistic documents at scale"""
    
    def __init__(self, total_documents: int = 5000):
        self.storage_path = Path("backend/data/storage")
        self.storage_path.mkdir(parents=True, exist_ok=True)
        
        self.total_documents = total_documents
        
        # Business domains with comprehensive templates
        self.document_domains = {
            "finance": {
                "weight": 0.20,
                "file_types": ["pdf", "xlsx", "docx"],
                "classifications": ["confidential", "internal", "private"],
                "templates": [
                    ("Financial Statement", "Comprehensive financial statements including balance sheets, income statements, and cash flow analysis for fiscal year {year}. Includes detailed revenue breakdown, expense categorization, asset valuation, liability assessment, and equity analysis with comparative data."),
                    ("Budget Proposal", "Detailed budget proposal for {dept} department covering Q{quarter} {year}. Includes operational expenses, capital expenditures, projected revenues, cost-benefit analysis, ROI calculations, and risk assessment with supporting documentation."),
                    ("Audit Report", "Independent audit report conducted by external auditors examining financial controls, compliance procedures, and operational efficiency for period ending {date}. Includes findings, recommendations, and management responses."),
                    ("Investment Analysis", "Comprehensive investment analysis for {project} initiative evaluating market opportunities, competitive landscape, financial projections, risk factors, and strategic alignment with organizational objectives."),
                    ("Tax Documentation", "Complete tax documentation for fiscal year {year} including federal and state returns, supporting schedules, deduction substantiation, estimated payments, and compliance certifications with audit trail."),
                    ("Cost Analysis Report", "Detailed cost analysis examining operational expenses, vendor contracts, resource allocation, efficiency metrics, and cost optimization opportunities for {dept} department."),
                    ("Revenue Forecast", "Multi-year revenue forecast incorporating market trends, historical performance, competitive analysis, and strategic initiatives. Includes sensitivity analysis and scenario planning."),
                    ("Expense Reconciliation", "Monthly expense reconciliation for {month} {year} with detailed transaction analysis, variance explanations, approval workflows, and supporting documentation."),
                ]
            },
            "legal": {
                "weight": 0.18,
                "file_types": ["pdf", "docx"],
                "classifications": ["confidential", "private"],
                "templates": [
                    ("Service Agreement", "Comprehensive service agreement between parties establishing scope of work, deliverables, timelines, payment terms, intellectual property rights, confidentiality obligations, liability limitations, dispute resolution procedures, and termination conditions."),
                    ("Non-Disclosure Agreement", "Bilateral non-disclosure agreement protecting confidential information, trade secrets, and proprietary data. Includes definition of confidential information, permitted uses, disclosure restrictions, return obligations, and breach remedies."),
                    ("Employment Contract", "Standard employment contract for {position} position detailing compensation structure, benefits package, duties and responsibilities, performance expectations, termination provisions, non-compete clauses, and confidentiality requirements."),
                    ("Compliance Policy", "Comprehensive compliance policy addressing regulatory requirements, industry standards, internal controls, reporting procedures, training obligations, and enforcement mechanisms for {area} compliance."),
                    ("Intellectual Property Agreement", "Intellectual property assignment and protection agreement covering inventions, patents, trademarks, copyrights, trade secrets, and proprietary information developed during employment or engagement."),
                    ("Vendor Contract", "Master vendor agreement establishing terms for procurement of goods and services including pricing, quality standards, delivery terms, warranties, indemnification, insurance requirements, and performance metrics."),
                    ("Partnership Agreement", "Strategic partnership agreement defining roles, responsibilities, profit sharing, decision-making authority, capital contributions, dispute resolution, and exit provisions between parties."),
                    ("License Agreement", "Software and technology license agreement granting usage rights, defining permitted uses, establishing restrictions, addressing support and maintenance, and specifying renewal terms."),
                ]
            },
            "operations": {
                "weight": 0.15,
                "file_types": ["pdf", "docx", "xlsx"],
                "classifications": ["internal", "public"],
                "templates": [
                    ("Standard Operating Procedure", "Detailed standard operating procedure for {process} covering step-by-step instructions, quality checkpoints, safety requirements, escalation procedures, required documentation, and continuous improvement mechanisms."),
                    ("Process Documentation", "Comprehensive process documentation for {system} including workflow diagrams, decision trees, input requirements, output specifications, exception handling, and integration points with related systems."),
                    ("Quality Control Manual", "Quality control and assurance manual establishing standards, testing procedures, acceptance criteria, non-conformance protocols, corrective actions, and continuous monitoring systems."),
                    ("Safety Protocol", "Workplace safety protocol addressing hazard identification, risk assessment, protective measures, emergency procedures, incident reporting, and regulatory compliance requirements."),
                    ("Maintenance Schedule", "Preventive maintenance schedule for {equipment} including inspection frequencies, service requirements, replacement criteria, downtime planning, and performance monitoring."),
                    ("Inventory Management", "Inventory management procedures covering receiving, storage, tracking, cycle counting, replenishment, obsolescence management, and reconciliation processes."),
                    ("Supply Chain Manual", "Supply chain operations manual addressing vendor management, procurement processes, logistics coordination, quality assurance, and performance metrics."),
                    ("Facility Operations", "Facility operations guide covering building management, utilities, security, maintenance, space planning, and emergency preparedness procedures."),
                ]
            },
            "hr": {
                "weight": 0.15,
                "file_types": ["pdf", "docx", "xlsx"],
                "classifications": ["internal", "confidential"],
                "templates": [
                    ("Employee Handbook", "Comprehensive employee handbook covering organizational policies, code of conduct, employment practices, benefits programs, leave policies, performance management, and disciplinary procedures."),
                    ("Performance Review", "Annual performance review for {position} evaluating goal achievement, competency demonstration, development progress, behavioral observations, and growth opportunities with ratings and recommendations."),
                    ("Recruitment Strategy", "Strategic recruitment plan for {dept} department identifying talent needs, sourcing strategies, selection criteria, onboarding processes, and success metrics."),
                    ("Training Program", "Comprehensive training program for {skill} development including learning objectives, curriculum design, delivery methods, assessment criteria, and certification requirements."),
                    ("Compensation Analysis", "Market-based compensation analysis examining salary structures, benefits packages, incentive programs, and total rewards strategy with competitive benchmarking data."),
                    ("Succession Planning", "Leadership succession planning document identifying critical roles, potential successors, development plans, knowledge transfer strategies, and transition timelines."),
                    ("Diversity & Inclusion Policy", "Diversity and inclusion policy establishing organizational commitment, goals, initiatives, accountability measures, and progress tracking for creating inclusive workplace."),
                    ("Employee Engagement Survey", "Comprehensive employee engagement survey results analyzing satisfaction, motivation, organizational culture, leadership effectiveness, and improvement opportunities."),
                ]
            },
            "marketing": {
                "weight": 0.12,
                "file_types": ["pdf", "pptx", "docx"],
                "classifications": ["internal", "public"],
                "templates": [
                    ("Marketing Strategy", "Integrated marketing strategy for {year} defining target markets, value propositions, channel strategies, budget allocation, campaign plans, and performance metrics."),
                    ("Brand Guidelines", "Comprehensive brand guidelines establishing visual identity, messaging frameworks, tone of voice, logo usage, color palettes, typography standards, and brand protection protocols."),
                    ("Campaign Analysis", "Multi-channel marketing campaign analysis for {campaign} examining reach, engagement, conversion rates, ROI, customer acquisition costs, and optimization recommendations."),
                    ("Market Research", "Primary and secondary market research report analyzing industry trends, competitive dynamics, customer preferences, market sizing, and strategic opportunities."),
                    ("Content Strategy", "Content marketing strategy defining audience personas, content pillars, editorial calendar, distribution channels, SEO optimization, and engagement metrics."),
                    ("Product Launch Plan", "Comprehensive product launch plan for {product} including market positioning, pricing strategy, promotional tactics, sales enablement, and success criteria."),
                    ("Customer Segmentation", "Data-driven customer segmentation analysis identifying distinct customer groups, behavioral patterns, preferences, lifetime value, and targeted engagement strategies."),
                    ("Digital Marketing Report", "Digital marketing performance report analyzing website traffic, social media engagement, email effectiveness, paid advertising ROI, and conversion optimization."),
                ]
            },
            "technology": {
                "weight": 0.10,
                "file_types": ["pdf", "docx", "txt", "json"],
                "classifications": ["internal", "private"],
                "templates": [
                    ("System Architecture", "Enterprise system architecture documentation defining infrastructure components, integration patterns, security layers, scalability design, disaster recovery, and technology standards."),
                    ("API Documentation", "RESTful API documentation covering endpoints, authentication methods, request/response formats, error handling, rate limiting, versioning strategy, and integration examples."),
                    ("Security Assessment", "Comprehensive cybersecurity assessment examining vulnerabilities, threat vectors, existing controls, compliance gaps, remediation priorities, and security roadmap."),
                    ("Technical Specification", "Detailed technical specification for {system} development including functional requirements, non-functional requirements, data models, interface designs, and acceptance criteria."),
                    ("Disaster Recovery Plan", "Enterprise disaster recovery and business continuity plan addressing risk scenarios, backup strategies, recovery procedures, RTO/RPO targets, and testing protocols."),
                    ("Database Design", "Database schema design documentation including entity relationships, data dictionary, indexing strategy, query optimization, backup procedures, and performance tuning."),
                    ("Integration Guide", "System integration guide for connecting {system} with enterprise applications including authentication, data mapping, error handling, and monitoring requirements."),
                    ("DevOps Runbook", "Operational runbook for DevOps team covering deployment procedures, monitoring dashboards, incident response, troubleshooting guides, and escalation protocols."),
                ]
            },
            "healthcare": {
                "weight": 0.05,
                "file_types": ["pdf", "docx"],
                "classifications": ["confidential", "private"],
                "templates": [
                    ("Clinical Protocol", "Evidence-based clinical protocol for {condition} management including diagnostic criteria, treatment algorithms, medication guidelines, monitoring requirements, and outcome measures."),
                    ("HIPAA Compliance", "HIPAA compliance documentation addressing privacy rules, security standards, breach notification procedures, business associate agreements, and audit readiness."),
                    ("Patient Care Manual", "Comprehensive patient care manual covering assessment procedures, care planning, intervention protocols, documentation requirements, and quality indicators."),
                    ("Medical Device Protocol", "Medical device operation and maintenance protocol including safety checks, calibration procedures, cleaning standards, troubleshooting guides, and regulatory compliance."),
                    ("Quality Improvement", "Quality improvement initiative documentation analyzing clinical outcomes, process metrics, patient satisfaction, adverse events, and continuous improvement actions."),
                ]
            },
            "real_estate": {
                "weight": 0.03,
                "file_types": ["pdf", "docx"],
                "classifications": ["internal", "confidential"],
                "templates": [
                    ("Property Appraisal", "Professional property appraisal report for {address} including comparative market analysis, property valuation, condition assessment, and valuation methodology."),
                    ("Lease Agreement", "Commercial lease agreement establishing rental terms, tenant obligations, landlord responsibilities, maintenance provisions, renewal options, and termination conditions."),
                    ("Property Management", "Property management procedures covering tenant relations, rent collection, maintenance coordination, vendor management, and financial reporting."),
                ]
            },
            "academic": {
                "weight": 0.02,
                "file_types": ["pdf", "docx"],
                "classifications": ["internal", "public"],
                "templates": [
                    ("Research Proposal", "Academic research proposal investigating {topic} with literature review, research methodology, expected contributions, timeline, and resource requirements."),
                    ("Course Curriculum", "Comprehensive course curriculum for {subject} including learning objectives, topic coverage, assessment methods, required readings, and grading criteria."),
                ]
            }
        }
        
        # User credentials for equitable distribution
        self.users_by_role = {}
        self.managers_by_dept = {}  # Track managers per department
        self.analysts_by_manager = {}  # Track which analysts belong to which manager
    
    async def generate_comprehensive_users(self, session: AsyncSession) -> Dict[str, User]:
        """
        Generate realistic business users with proper organizational hierarchy.
        
        Structure:
        - Admin: Top-level executives (2-3 total)
        - Manager: Department heads (1-2 per department, 10-14 total)
        - Analyst: Team members (5-8 per Manager, 60-80 total)
        
        Hierarchy: Admin → Manager → Analyst (with manager_id relationships)
        """
        logger.info("👥 Creating enterprise user base with organizational hierarchy...")
        
        # Department structure with role distribution
        departments = [
            {"name": "Finance", "managers": 2, "analysts_per_manager": 6},
            {"name": "Legal", "managers": 2, "analysts_per_manager": 5},
            {"name": "Operations", "managers": 2, "analysts_per_manager": 7},
            {"name": "HR", "managers": 1, "analysts_per_manager": 8},
            {"name": "Marketing", "managers": 2, "analysts_per_manager": 6},
            {"name": "Technology", "managers": 2, "analysts_per_manager": 8},
            {"name": "Executive", "managers": 1, "analysts_per_manager": 5},
        ]
        
        created_users = {}
        
        # Phase 1: Create Admins (Executive leadership)
        logger.info("  Creating Admin users...")
        admin_titles = ["Chief Executive Officer", "Chief Technology Officer", "Chief Operating Officer"]
        for idx, title in enumerate(admin_titles):
            name = f"{fake.first_name()} {fake.last_name()}"
            email = f"{name.lower().replace(' ', '.')}@enterprise.indoc.local"
            username = f"admin_{name.lower().replace(' ', '_')}"
            
            # Check if user exists
            result = await session.execute(select(User).where(User.email == email))
            existing_user = result.scalar_one_or_none()
            
            if existing_user:
                user = existing_user
            else:
                import secrets
                user = User(
                    email=email,
                    username=username,
                    full_name=name,
                    hashed_password=get_password_hash(secrets.token_urlsafe(16)),
                    role=UserRole.ADMIN,
                    department=title,
                    is_active=True,
                    is_verified=True,
                    manager_id=None  # Admins have no manager
                )
                session.add(user)
                await session.flush()
                logger.debug(f"  Created Admin: {name} ({title})")
            
            key = f"Admin_Executive_{idx}"
            created_users[key] = {
                "user": user,
                "department": "Executive",
                "role": UserRole.ADMIN
            }
            
            if UserRole.ADMIN not in self.users_by_role:
                self.users_by_role[UserRole.ADMIN] = []
            self.users_by_role[UserRole.ADMIN].append(user)
        
        # Phase 2: Create Managers per department
        logger.info("  Creating Manager users...")
        for dept_config in departments:
            dept_name = dept_config["name"]
            self.managers_by_dept[dept_name] = []
            
            for mgr_idx in range(dept_config["managers"]):
                name = f"{fake.first_name()} {fake.last_name()}"
                email = f"{name.lower().replace(' ', '.')}@enterprise.indoc.local"
                username = f"mgr_{name.lower().replace(' ', '_')}"
                
                # Check if user exists
                result = await session.execute(select(User).where(User.email == email))
                existing_user = result.scalar_one_or_none()
                
                if existing_user:
                    user = existing_user
                else:
                    import secrets
                    user = User(
                        email=email,
                        username=username,
                        full_name=name,
                        hashed_password=get_password_hash(secrets.token_urlsafe(16)),
                        role=UserRole.MANAGER,
                        department=dept_name,
                        is_active=True,
                        is_verified=True,
                        manager_id=None  # Managers report to Admins (not tracked in this simple hierarchy)
                    )
                    session.add(user)
                    await session.flush()
                    logger.debug(f"  Created Manager: {name} ({dept_name})")
                
                key = f"Manager_{dept_name}_{mgr_idx}"
                created_users[key] = {
                    "user": user,
                    "department": dept_name,
                    "role": UserRole.MANAGER
                }
                
                if UserRole.MANAGER not in self.users_by_role:
                    self.users_by_role[UserRole.MANAGER] = []
                self.users_by_role[UserRole.MANAGER].append(user)
                self.managers_by_dept[dept_name].append(user)
                
                # Initialize analysts list for this manager
                self.analysts_by_manager[user.id] = []
        
        # Phase 3: Create Analysts and assign to Managers
        logger.info("  Creating Analyst users and building hierarchy...")
        for dept_config in departments:
            dept_name = dept_config["name"]
            managers_in_dept = self.managers_by_dept[dept_name]
            
            # Calculate total analysts for this department
            total_analysts = dept_config["managers"] * dept_config["analysts_per_manager"]
            
            for analyst_idx in range(total_analysts):
                name = f"{fake.first_name()} {fake.last_name()}"
                email = f"{name.lower().replace(' ', '.')}@enterprise.indoc.local"
                username = f"analyst_{name.lower().replace(' ', '_')}"
                
                # Assign to manager in round-robin fashion for even distribution
                manager = managers_in_dept[analyst_idx % len(managers_in_dept)]
                
                # Check if user exists
                result = await session.execute(select(User).where(User.email == email))
                existing_user = result.scalar_one_or_none()
                
                if existing_user:
                    user = existing_user
                    # Update manager_id if not set
                    if not user.manager_id:
                        user.manager_id = manager.id
                else:
                    import secrets
                    user = User(
                        email=email,
                        username=username,
                        full_name=name,
                        hashed_password=get_password_hash(secrets.token_urlsafe(16)),
                        role=UserRole.ANALYST,
                        department=dept_name,
                        is_active=True,
                        is_verified=True,
                        manager_id=manager.id  # KEY: Link analyst to their manager
                    )
                    session.add(user)
                    await session.flush()
                    logger.debug(f"  Created Analyst: {name} ({dept_name}, reports to {manager.full_name})")
                
                key = f"Analyst_{dept_name}_{analyst_idx}"
                created_users[key] = {
                    "user": user,
                    "department": dept_name,
                    "role": UserRole.ANALYST,
                    "manager": manager
                }
                
                if UserRole.ANALYST not in self.users_by_role:
                    self.users_by_role[UserRole.ANALYST] = []
                self.users_by_role[UserRole.ANALYST].append(user)
                self.analysts_by_manager[manager.id].append(user)
        
        await session.commit()
        
        # Print hierarchy summary
        logger.info(f"\n  ✅ Created {len(created_users)} enterprise users:")
        logger.info(f"     • {len(self.users_by_role.get(UserRole.ADMIN, []))} Admins")
        logger.info(f"     • {len(self.users_by_role.get(UserRole.MANAGER, []))} Managers")
        logger.info(f"     • {len(self.users_by_role.get(UserRole.ANALYST, []))} Analysts")
        logger.info(f"  📊 Hierarchy established:")
        for manager_id, analysts in self.analysts_by_manager.items():
            if analysts:
                manager = next((u for u in self.users_by_role[UserRole.MANAGER] if u.id == manager_id), None)
                if manager:
                    logger.info(f"     • {manager.full_name} → {len(analysts)} analysts")
        
        return created_users
    
    async def generate_realistic_document_content(
        self, 
        template: tuple, 
        domain: str, 
        file_type: str,
        index: int
    ) -> tuple:
        """Generate realistic document content with comprehensive details"""
        
        title_template, desc_template = template
        
        # Generate contextual variables
        year = random.randint(2022, 2024)
        quarter = random.randint(1, 4)
        month = fake.month_name()
        dept = random.choice(["Finance", "Operations", "Technology", "Marketing", "HR", "Legal"])
        project = fake.catch_phrase()
        process = fake.bs()
        system = f"{fake.word().capitalize()} {random.choice(['System', 'Platform', 'Application'])}"
        position = fake.job()
        address = fake.address().replace('\n', ', ')
        skill = fake.job()
        campaign = f"{fake.color_name().capitalize()} {fake.word().capitalize()} Campaign"
        product = fake.catch_phrase()
        equipment = f"{fake.word().capitalize()} {random.choice(['Equipment', 'Machinery', 'Device'])}"
        condition = f"{fake.word().capitalize()} {random.choice(['Condition', 'Disorder', 'Syndrome'])}"
        area = random.choice(["Financial", "Data Privacy", "Safety", "Environmental", "Quality"])
        topic = fake.catch_phrase()
        subject = f"{fake.word().capitalize()} {random.choice(['Studies', 'Science', 'Management'])}"
        date = fake.date_between(start_date='-2y', end_date='today').strftime("%B %Y")
        
        # Format title and description
        title = title_template.format(
            year=year, quarter=quarter, month=month, dept=dept, project=project,
            process=process, system=system, position=position, address=address,
            skill=skill, campaign=campaign, product=product, equipment=equipment,
            condition=condition, area=area, topic=topic, subject=subject, date=date
        )
        
        description = desc_template.format(
            year=year, quarter=quarter, month=month, dept=dept, project=project,
            process=process, system=system, position=position, address=address,
            skill=skill, campaign=campaign, product=product, equipment=equipment,
            condition=condition, area=area, topic=topic, subject=subject, date=date
        )
        
        # Generate comprehensive content
        content_sections = [
            f"DOCUMENT ID: DOC-{domain.upper()}-{index:06d}",
            f"CLASSIFICATION: {domain.capitalize()} - {title}",
            f"GENERATED: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            "",
            "EXECUTIVE SUMMARY",
            "=" * 80,
            description,
            "",
            "DETAILED CONTENT",
            "=" * 80,
        ]
        
        # Add domain-specific content
        for section_num in range(random.randint(5, 10)):
            content_sections.extend([
                "",
                f"{section_num + 1}. {fake.catch_phrase().upper()}",
                "-" * 80,
                *[fake.paragraph(nb_sentences=random.randint(3, 6)) for _ in range(random.randint(2, 4))],
            ])
        
        # Add appendix
        content_sections.extend([
            "",
            "APPENDIX",
            "=" * 80,
            *[f"• {fake.sentence()}" for _ in range(random.randint(5, 10))],
            "",
            f"Document Version: 1.{random.randint(0, 9)}",
            f"Last Updated: {fake.date_between(start_date='-30d', end_date='today')}",
            f"Next Review Date: {fake.date_between(start_date='today', end_date='+1y')}",
            "",
            "END OF DOCUMENT"
        ])
        
        full_content = "\n".join(content_sections)
        
        # Generate tags
        tags = [
            domain,
            str(year),
            dept.lower(),
            random.choice(["strategic", "operational", "tactical"]),
            random.choice(["quarterly", "annual", "monthly", "ongoing"])
        ]
        
        return title, description, full_content, tags
    
    def create_file_bytes(self, content: str, file_type: str) -> bytes:
        """Create actual file bytes for different file types"""
        try:
            if file_type == "docx":
                try:
                    from docx import Document as DocxDocument
                    doc = DocxDocument()
                    for paragraph in content.split("\n"):
                        doc.add_paragraph(paragraph)
                    buffer = BytesIO()
                    doc.save(buffer)
                    return buffer.getvalue()
                except ImportError:
                    pass
            
            elif file_type == "xlsx":
                try:
                    from openpyxl import Workbook
                    wb = Workbook()
                    ws = wb.active
                    ws.title = "Data"
                    for i, line in enumerate(content.split("\n")[:1000], start=1):
                        ws.cell(row=i, column=1, value=line)
                    buffer = BytesIO()
                    wb.save(buffer)
                    return buffer.getvalue()
                except ImportError:
                    pass
            
            elif file_type == "pptx":
                try:
                    from pptx import Presentation
                    from pptx.util import Pt
                    prs = Presentation()
                    
                    # Title slide
                    slide_layout = prs.slide_layouts[0]
                    slide = prs.slides.add_slide(slide_layout)
                    title = slide.shapes.title
                    title.text = content.split("\n")[0][:100]
                    
                    # Content slides
                    for section in content.split("\n\n")[:10]:
                        if section.strip():
                            slide_layout = prs.slide_layouts[1]
                            slide = prs.slides.add_slide(slide_layout)
                            title = slide.shapes.title
                            body = slide.placeholders[1]
                            title.text = section[:50]
                            body.text = section[:500]
                    
                    buffer = BytesIO()
                    prs.save(buffer)
                    return buffer.getvalue()
                except ImportError:
                    pass
            
            elif file_type == "pdf":
                try:
                    from reportlab.lib.pagesizes import letter
                    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                    from reportlab.lib.styles import getSampleStyleSheet
                    
                    buffer = BytesIO()
                    doc = SimpleDocTemplate(buffer, pagesize=letter)
                    styles = getSampleStyleSheet()
                    story = []
                    
                    for line in content.split("\n")[:500]:
                        if line.strip():
                            story.append(Paragraph(line, styles['Normal']))
                            story.append(Spacer(1, 12))
                    
                    doc.build(story)
                    return buffer.getvalue()
                except ImportError:
                    pass
            
            elif file_type == "json":
                import json
                data = {
                    "document": {
                        "content": content[:1000],
                        "metadata": {
                            "generated": datetime.now().isoformat(),
                            "format": "json",
                            "lines": len(content.split("\n"))
                        }
                    }
                }
                return json.dumps(data, indent=2).encode('utf-8')
        
        except Exception as e:
            logger.debug(f"File generation error for {file_type}: {e}, falling back to text")
        
        # Fallback to plain text
        return content.encode('utf-8')
    
    def _select_smart_uploader(self, domain: str, doc_counter: int) -> User:
        """
        Intelligently select document uploader based on:
        1. Domain-to-department mapping
        2. Role distribution (70% Analyst, 20% Manager, 10% Admin)
        3. Even distribution across users
        """
        # Map domains to departments
        domain_dept_mapping = {
            "finance": "Finance",
            "legal": "Legal",
            "operations": "Operations",
            "hr": "HR",
            "marketing": "Marketing",
            "technology": "Technology",
            "healthcare": "Operations",  # Cross-functional
            "real_estate": "Operations",
            "academic": "Executive"
        }
        
        target_dept = domain_dept_mapping.get(domain, "Operations")
        
        # Determine role based on distribution (70% Analyst, 20% Manager, 10% Admin)
        roll = (doc_counter * 37) % 100  # Pseudo-random but deterministic
        
        if roll < 70:
            # 70% - Assign to Analyst
            target_role = UserRole.ANALYST
            # Get analysts from target department
            candidates = [
                u for u in self.users_by_role.get(UserRole.ANALYST, [])
                if getattr(u, 'department', None) == target_dept
            ]
            if not candidates:
                # Fallback to any analyst
                candidates = self.users_by_role.get(UserRole.ANALYST, [])
        
        elif roll < 90:
            # 20% - Assign to Manager
            target_role = UserRole.MANAGER
            # Get managers from target department
            candidates = [
                u for u in self.users_by_role.get(UserRole.MANAGER, [])
                if getattr(u, 'department', None) == target_dept
            ]
            if not candidates:
                # Fallback to any manager
                candidates = self.users_by_role.get(UserRole.MANAGER, [])
        
        else:
            # 10% - Assign to Admin
            target_role = UserRole.ADMIN
            candidates = self.users_by_role.get(UserRole.ADMIN, [])
        
        # If no candidates found, fallback to any user
        if not candidates:
            all_users = []
            for role_users in self.users_by_role.values():
                all_users.extend(role_users)
            candidates = all_users
        
        # Round-robin selection within candidates for even distribution
        if candidates:
            return candidates[doc_counter % len(candidates)]
        
        # Last resort - return first user
        return list(self.users_by_role.values())[0][0]
    
    async def generate_enterprise_documents(
        self, 
        session: AsyncSession, 
        users: Dict[str, Any]
    ) -> List[Document]:
        """Generate thousands of enterprise documents"""
        logger.info(f"📄 Generating {self.total_documents} enterprise documents...")
        
        documents = []
        commit_batch = 100
        
        # Calculate documents per domain based on weights
        domain_counts = {}
        remaining = self.total_documents
        
        for domain, config in self.document_domains.items():
            count = int(self.total_documents * config["weight"])
            domain_counts[domain] = count
            remaining -= count
        
        # Distribute remaining to largest domain
        largest_domain = max(self.document_domains.keys(), key=lambda k: self.document_domains[k]["weight"])
        domain_counts[largest_domain] += remaining
        
        doc_counter = 0
        
        for domain, count in domain_counts.items():
            config = self.document_domains[domain]
            templates = config["templates"]
            
            logger.info(f"  Generating {count} {domain} documents...")
            
            for i in range(count):
                doc_counter += 1
                
                # Select template and file type
                template = random.choice(templates)
                file_type = random.choice(config["file_types"])
                classification = random.choice(config["classifications"])
                
                # Generate content
                base_title, description, content, tags = await self.generate_realistic_document_content(
                    template, domain, file_type, doc_counter
                )
                
                # Make title unique by adding counter
                title = f"{base_title} - {doc_counter:06d}"
                
                # Create file
                file_bytes = self.create_file_bytes(content, file_type)
                file_hash = hashlib.sha256(file_bytes).hexdigest()
                
                # Save to local storage (primary)
                storage_file = self.storage_path / f"{file_hash}.{file_type}"
                storage_file.write_bytes(file_bytes)
                
                # Dual-write to S3 (if configured)
                s3_success = False
                try:
                    from app.services.storage.factory import get_primary_storage
                    from app.services.storage.base import build_object_key
                    from app.core.config import settings
                    
                    s3_client = get_primary_storage()
                    if s3_client:
                        # Build S3 key using configured prefix
                        object_key = build_object_key(
                            tenant_id=None,  # Seed data has no tenant
                            file_hash=file_hash,
                            file_extension=file_type,
                            prefix=settings.S3_PREFIX  # Use configured prefix: 'file-storage'
                        )
                        
                        # Upload to S3
                        s3_client.put_bytes(object_key, file_bytes, content_type=f'application/{file_type}')
                        s3_success = True
                        logger.debug(f"☁️  Uploaded {title[:30]} to S3: {object_key}")
                except Exception as e:
                    logger.debug(f"S3 upload skipped for {title[:30]}: {e}")
                
                # Smart uploader selection: 70% Analyst, 20% Manager, 10% Admin
                # Matched to appropriate department based on document domain
                uploader = self._select_smart_uploader(domain, doc_counter)
                
                # Create document record
                custom_meta = {
                    "domain": domain,
                    "generated": datetime.now().isoformat(),
                    "enterprise_seed": True,
                    "version": "1.0"
                }
                
                # Add S3 key if successfully uploaded
                if s3_success:
                    custom_meta["object_storage_key"] = object_key
                    custom_meta["s3_uploaded"] = True
                
                doc_kwargs = {
                    "uuid": uuid.uuid4(),
                    "filename": f"{title[:50].replace(' ', '_')}_{doc_counter}.{file_type}",
                    "file_type": file_type,
                    "file_size": len(file_bytes),
                    "file_hash": file_hash,
                    "storage_path": str(storage_file),
                    "status": "pending",  # Will be processed by Celery
                    "virus_scan_status": "clean",
                    "title": title,
                    "description": description,
                    "tags": tags,
                    "full_text": content,
                    "language": "en",
                    "access_level": classification,
                    "uploaded_by": uploader.id,
                    "custom_metadata": custom_meta
                }
                
                # Add classification if the field exists
                classification_enum = self._map_classification(classification)
                if classification_enum is not None:
                    doc_kwargs["classification"] = classification_enum
                
                document = Document(**doc_kwargs)
                
                session.add(document)
                documents.append(document)
                
                # Commit in batches
                if len(documents) % commit_batch == 0:
                    await session.flush()
                    await session.commit()
                    logger.info(f"    Progress: {len(documents)}/{self.total_documents} documents created")
        
        await session.flush()
        await session.commit()
        logger.info(f"✅ Generated {len(documents)} enterprise documents")
        
        return documents
    
    def _map_classification(self, access_level: str):
        """Map access level to DocumentClassification enum"""
        try:
            from app.models.classification import DocumentClassification
            
            mapping = {
                "public": DocumentClassification.PUBLIC,
                "internal": DocumentClassification.INTERNAL,
                "confidential": DocumentClassification.CONFIDENTIAL,
                "private": DocumentClassification.RESTRICTED
            }
            return mapping.get(access_level, DocumentClassification.INTERNAL)
        except ImportError:
            # Classification module doesn't exist, return None
            return None
    
    async def generate_document_permissions(
        self, 
        session: AsyncSession, 
        documents: List[Document]
    ):
        """
        Generate explicit document permissions for audit trail and access control.
        
        Rules:
        - Document owner gets full permissions (read, write, share, delete)
        - Owner's manager gets read permission (for hierarchy access)
        - All Admins get read permission (for oversight)
        - Random share permissions for collaboration scenarios
        """
        logger.info("🔐 Generating document permissions for access control and audit trail...")
        
        try:
            from app.models.document_permission import DocumentPermission
        except ImportError:
            logger.warning("⚠️  DocumentPermission model not found, skipping permission generation")
            return
        
        permissions_created = 0
        batch_size = 100
        
        for doc in documents:
            # Get document owner
            owner_id = doc.uploaded_by
            owner = await session.get(User, owner_id)
            
            if not owner:
                continue
            
            # Track granted permissions to avoid duplicates (document_id, user_id, permission_type must be unique)
            granted_permissions = set()
            
            # 1. Grant full permissions to document owner
            for perm_type in ['read', 'write', 'share', 'delete']:
                key = (doc.id, owner.id, perm_type)
                if key not in granted_permissions:
                    permission = DocumentPermission(
                        document_id=doc.id,
                        user_id=owner.id,
                        permission_type=perm_type,
                        granted_by=owner.id,  # Self-granted
                        reason=f"Document owner - automatic grant"
                    )
                    session.add(permission)
                    permissions_created += 1
                    granted_permissions.add(key)
            
            # 2. Grant read permission to owner's manager (if exists)
            if owner.manager_id:
                manager = await session.get(User, owner.manager_id)
                if manager:
                    key = (doc.id, manager.id, 'read')
                    if key not in granted_permissions:
                        permission = DocumentPermission(
                            document_id=doc.id,
                            user_id=manager.id,
                            permission_type='read',
                            granted_by=owner.id,
                            reason="Manager hierarchy access"
                        )
                        session.add(permission)
                        permissions_created += 1
                        granted_permissions.add(key)
            
            # 3. Grant read permission to all Admins
            for admin in self.users_by_role.get(UserRole.ADMIN, []):
                key = (doc.id, admin.id, 'read')
                if key not in granted_permissions:
                    permission = DocumentPermission(
                        document_id=doc.id,
                        user_id=admin.id,
                        permission_type='read',
                        granted_by=owner.id,
                        reason="Admin oversight access"
                    )
                    session.add(permission)
                    permissions_created += 1
                    granted_permissions.add(key)
            
            # 4. Add some collaborative permissions (10% of documents)
            # Simulate real-world sharing scenarios
            if random.random() < 0.10 and owner.role == UserRole.ANALYST:
                # Pick 1-3 random colleagues from same department
                same_dept_users = [
                    u for u in self.users_by_role.get(UserRole.ANALYST, [])
                    if getattr(u, 'department', None) == getattr(owner, 'department', None)
                    and u.id != owner.id
                ]
                
                if same_dept_users:  # Only proceed if there are potential collaborators
                    num_collaborators = random.randint(1, min(3, len(same_dept_users)))
                    collaborators = random.sample(same_dept_users, num_collaborators)
                    
                    for collab in collaborators:
                        # Grant read permission to collaborator
                        key = (doc.id, collab.id, 'read')
                        if key not in granted_permissions:
                            permission = DocumentPermission(
                                document_id=doc.id,
                                user_id=collab.id,
                                permission_type='read',
                                granted_by=owner.id,
                                reason="Shared for collaboration"
                            )
                            session.add(permission)
                            permissions_created += 1
                            granted_permissions.add(key)
            
            # Commit in batches for performance
            if permissions_created % batch_size == 0:
                await session.flush()
                logger.debug(f"  Flushed {permissions_created} permissions...")
        
        await session.commit()
        logger.info(f"✅ Created {permissions_created} document permissions")
        logger.info(f"   • Owner permissions: {len(documents) * 4}")
        logger.info(f"   • Manager hierarchy permissions: ~{len([d for d in documents if (await session.get(User, d.uploaded_by)).manager_id])}")
        logger.info(f"   • Admin oversight permissions: {len(documents) * len(self.users_by_role.get(UserRole.ADMIN, []))}")
        logger.info(f"   • Collaborative shares: ~{int(len(documents) * 0.10 * 2)}")
    
    async def trigger_document_processing(self, documents: List[Document]):
        """Trigger Celery tasks to process and index documents"""
        logger.info("🔄 Triggering document processing for Elasticsearch and Qdrant indexing...")
        
        try:
            from app.tasks.document import process_document
            
            # Process in batches to avoid overwhelming Celery
            batch_size = 50
            for i in range(0, len(documents), batch_size):
                batch = documents[i:i+batch_size]
                for doc in batch:
                    try:
                        process_document.delay(str(doc.uuid))
                    except Exception as e:
                        logger.warning(f"Failed to queue document {doc.uuid}: {e}")
                
                logger.info(f"  Queued {min(i+batch_size, len(documents))}/{len(documents)} documents for processing")
            
            logger.info("✅ Document processing tasks queued successfully")
            logger.info("   Note: Processing will continue in background via Celery workers")
        
        except Exception as e:
            logger.error(f"❌ Failed to trigger document processing: {e}")
            logger.info("   Documents created but not indexed. Run init_search_indices.py to index manually")
    
    async def generate_audit_trail(self, session: AsyncSession, users: Dict[str, Any], documents: List[Document]):
        """Generate realistic audit trail"""
        logger.info("📊 Generating audit trail...")
        
        actions = ["upload", "view", "download", "search", "update", "delete"]
        resource_types = ["document", "user", "search", "audit"]
        
        # Generate realistic audit entries
        for _ in range(min(1000, len(documents) * 2)):
            user_data = random.choice(list(users.values()))
            user = user_data["user"]
            
            action = random.choice(actions)
            resource_type = random.choice(resource_types)
            
            if resource_type == "document" and documents:
                resource_id = str(random.choice(documents).id)
            elif resource_type == "search":
                resource_id = fake.word()
            else:
                resource_id = str(user.id)
            
            audit_log = AuditLog(
                user_id=user.id,
                user_email=user.email,
                user_role=user.role.value,
                action=action,
                resource_type=resource_type,
                resource_id=resource_id,
                ip_address=fake.ipv4_private(),
                user_agent=fake.user_agent(),
                metadata={
                    "department": user_data["department"],
                    "generated": datetime.now().isoformat()
                }
            )
            session.add(audit_log)
        
        await session.commit()
        logger.info("✅ Audit trail generated")
    
    async def print_summary(self, session: AsyncSession, users: Dict[str, Any]):
        """Print generation summary"""
        logger.info("\n" + "=" * 80)
        logger.info("📋 ENTERPRISE SEED DATA GENERATION SUMMARY")
        logger.info("=" * 80)
        
        # User summary
        logger.info("\n👥 USER SUMMARY:")
        for role in UserRole:
            role_users = [u for u in users.values() if u["role"] == role]
            if role_users:
                logger.info(f"  {role.value.capitalize():15} : {len(role_users):3} users")
        
        # Document summary
        result = await session.execute(select(Document))
        all_docs = result.scalars().all()
        
        logger.info(f"\n📄 DOCUMENT SUMMARY: {len(all_docs)} total documents")
        
        # By domain
        domains = {}
        for doc in all_docs:
            if doc.custom_metadata and "domain" in doc.custom_metadata:
                domain = doc.custom_metadata["domain"]
                domains[domain] = domains.get(domain, 0) + 1
        
        logger.info("\n  By Domain:")
        for domain, count in sorted(domains.items(), key=lambda x: x[1], reverse=True):
            logger.info(f"    {domain.capitalize():15} : {count:4} documents")
        
        # By file type
        file_types = {}
        for doc in all_docs:
            file_types[doc.file_type] = file_types.get(doc.file_type, 0) + 1
        
        logger.info("\n  By File Type:")
        for ft, count in sorted(file_types.items(), key=lambda x: x[1], reverse=True):
            logger.info(f"    {ft.upper():15} : {count:4} documents")
        
        # By classification
        classifications = {}
        for doc in all_docs:
            classifications[doc.access_level] = classifications.get(doc.access_level, 0) + 1
        
        logger.info("\n  By Classification:")
        for cls, count in sorted(classifications.items(), key=lambda x: x[1], reverse=True):
            logger.info(f"    {cls.capitalize():15} : {count:4} documents")
        
        logger.info("\n🔐 SAMPLE LOGIN CREDENTIALS:")
        logger.info("  Username: Any user email")
        logger.info("  Password: Enterprise2024!")
        
        logger.info("\n" + "=" * 80)
        logger.info("✅ Enterprise seed data generation complete!")
        logger.info("=" * 80 + "\n")
    
    async def run(self, clean_existing: bool = False):
        """Run the enterprise seed generation"""
        logger.info("🚀 Starting Enterprise Seed Data Generation")
        logger.info(f"   Target: {self.total_documents} documents")
        
        async with AsyncSessionLocal() as session:
            if clean_existing:
                logger.warning("🧹 Cleaning existing data...")
                # Delete in correct order to respect foreign key constraints
                await session.execute(delete(AuditLog))
                
                # Delete document permissions before documents
                try:
                    from app.models.document_permission import DocumentPermission
                    await session.execute(delete(DocumentPermission))
                    logger.info("  Cleaned document permissions")
                except ImportError:
                    pass
                
                # Import Conversation and Message if they exist
                try:
                    from app.models.conversation import Conversation, Message
                    await session.execute(delete(Message))
                    await session.execute(delete(Conversation))
                except ImportError:
                    pass
                
                await session.execute(delete(Document))
                await session.execute(delete(User).where(User.email.like('%@enterprise.indoc.local')))
                await session.commit()
                logger.info("✅ Existing data cleaned")
            
            # Generate users
            users = await self.generate_comprehensive_users(session)
            await session.commit()
            
            # Generate documents
            documents = await self.generate_enterprise_documents(session, users)
            
            # Generate document permissions for access control and audit
            await self.generate_document_permissions(session, documents)
            
            # Generate audit trail
            await self.generate_audit_trail(session, users, documents)
            
            # Print summary
            await self.print_summary(session, users)
            
            # Trigger processing (ES/Qdrant indexing)
            await self.trigger_document_processing(documents)


async def main():
    """Main entry point"""
    import argparse
    
    parser = argparse.ArgumentParser(description='Enterprise Seed Data Generator')
    parser.add_argument('--total', type=int, default=5000, help='Total documents to generate (default: 5000)')
    parser.add_argument('--clean', action='store_true', help='Clean existing data first')
    parser.add_argument('--yes', action='store_true', help='Auto-confirm prompts')
    
    args = parser.parse_args()
    
    # Environment variable override
    if os.getenv("INDOC_YES") == "1":
        args.yes = True
    
    if args.clean and not args.yes:
        logger.warning("⚠️  This will DELETE all existing documents and enterprise users!")
        response = input("Continue? (y/N): ")
        if response.lower() != 'y':
            logger.info("Cancelled")
            return
    
    generator = EnterpriseDocumentGenerator(total_documents=args.total)
    await generator.run(clean_existing=args.clean)


if __name__ == "__main__":
    asyncio.run(main())

